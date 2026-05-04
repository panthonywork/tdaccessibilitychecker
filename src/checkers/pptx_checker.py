from src.standards.matterhorn import get_checkpoints_for_filetype


def check_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    try:
        from pptx import Presentation
        from io import BytesIO
    except ImportError as e:
        return [_error_result(str(e))]

    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as e:
        return [_error_result(f"Could not open PowerPoint file: {e}")]

    checkpoints = get_checkpoints_for_filetype("pptx")
    results = []
    for cp in checkpoints:
        result = _run_checkpoint(cp, prs)
        results.append(result)
    return results


def _run_checkpoint(cp: dict, prs) -> dict:
    cp_id = cp["id"]
    base = {
        "id": cp_id,
        "group_id": cp["group_id"],
        "group_name": cp["group_name"],
        "description": cp["description"],
        "plain_english": cp["plain_english"],
        "fix_hint": cp["fix_hint"],
    }

    try:
        if cp_id == "06-001":
            status = _check_title(prs)
        elif cp_id == "07-001":
            status = _check_language(prs)
        elif cp_id == "17-001":
            status = _check_image_alt_text(prs)
        elif cp_id == "13-001":
            status = _check_slide_titles(prs)
        elif cp_id == "28-002":
            status = _check_link_text(prs)
        else:
            status = "na"
    except Exception:
        status = "na"

    return {**base, "status": status}


def _check_title(prs) -> str:
    props = prs.core_properties
    title = getattr(props, "title", "") or ""
    return "pass" if title.strip() else "fail"


def _check_language(prs) -> str:
    from pptx.oxml.ns import qn
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        lang = run._r.get(qn("lang"), None)
                        if lang:
                            return "pass"
    return "fail"


def _check_image_alt_text(prs) -> str:
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt_text = shape.name or ""
                # python-pptx exposes alt text via the nvSpPr descr attribute
                try:
                    descr = shape._element.nvPicPr.cNvPr.get("descr", "").strip()
                    if not descr:
                        return "fail"
                except Exception:
                    return "fail"
    return "pass"


def _check_slide_titles(prs) -> str:
    for slide in prs.slides:
        title_shape = slide.shapes.title
        if title_shape is None or not title_shape.text.strip():
            return "fail"
    return "pass"


def _check_link_text(prs) -> str:
    bad_phrases = {"click here", "here", "read more", "learn more", "more", "link"}
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.hyperlink and run.hyperlink.address:
                            if run.text.strip().lower() in bad_phrases:
                                return "fail"
    return "pass"


def _error_result(message: str) -> dict:
    return {
        "id": "ERR",
        "group_id": "00",
        "group_name": "Error",
        "description": message,
        "plain_english": message,
        "fix_hint": "Ensure the file is a valid .pptx file.",
        "status": "fail",
    }
